import os

import fitz
from docx import Document as DocxDocument
from pptx import Presentation


class DocService:
    def extract_text_from_file(self, file_path: str) -> str:
        _, extension = os.path.splitext(file_path)
        extension = extension.lower()

        if extension == ".docx":
            return self._extract_text_from_docx(file_path)
        elif extension == ".pdf":
            return self._extract_text_from_pdf(file_path)
        elif extension == ".pptx":
            return self._extract_text_from_pptx(file_path)
        else:
            raise ValueError(f"Unsupported file type: {extension}")

    @staticmethod
    def _extract_text_from_docx(file_path: str) -> str:
        doc = DocxDocument(file_path)
        return "\n".join([paragraph.text for paragraph in doc.paragraphs])

    @staticmethod
    def _extract_text_from_pdf(file_path: str) -> str:
        text = ""
        with fitz.open(file_path) as pdf:
            for page in pdf:
                text += page.get_text()
        return text

    @staticmethod
    def _extract_text_from_pptx(file_path: str) -> str:
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
